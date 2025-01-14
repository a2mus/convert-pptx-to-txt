from PyQt6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QLabel, QPushButton, QTextEdit, QFileDialog, QProgressBar,
    QMessageBox, QMenuBar, QMenu, QRadioButton, QButtonGroup
)
from PyQt6.QtCore import Qt
from PyQt6.QtGui import QAction
from pptx import Presentation
from spire.presentation import Presentation as SpirePresentation
import sys
import markdown
from markitdown import MarkItDown

class PPTXExtractorApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("PPTX Text Extractor")
        self.setGeometry(100, 100, 800, 600)
        self.init_ui()
        
    def init_ui(self):
        # Create main widget and layout
        main_widget = QWidget()
        self.setCentralWidget(main_widget)
        layout = QVBoxLayout(main_widget)
        
        # Create menu bar
        self.create_menu_bar()
        
        # Conversion type selection
        self.conversion_type = QButtonGroup()
        conversion_layout = QHBoxLayout()
        
        pptx_radio = QRadioButton("PPTX Conversion")
        pptx_radio.setChecked(True)
        self.conversion_type.addButton(pptx_radio, 0)
        
        markitdown_radio = QRadioButton("MarkItDown Conversion")
        self.conversion_type.addButton(markitdown_radio, 1)
        
        spire_radio = QRadioButton("Spire Conversion")
        self.conversion_type.addButton(spire_radio, 2)
        
        conversion_layout.addWidget(pptx_radio)
        conversion_layout.addWidget(markitdown_radio)
        conversion_layout.addWidget(spire_radio)
        layout.addLayout(conversion_layout)

        # File selection area
        file_layout = QHBoxLayout()
        self.file_label = QLabel("No file selected")
        self.file_label.setStyleSheet("font-weight: bold;")
        browse_btn = QPushButton("Browse...")
        browse_btn.clicked.connect(self.select_file)
        file_layout.addWidget(self.file_label)
        file_layout.addWidget(browse_btn)
        layout.addLayout(file_layout)
        
        # Progress bar
        self.progress = QProgressBar()
        self.progress.setVisible(False)
        layout.addWidget(self.progress)
        
        # Text display area
        self.text_area = QTextEdit()
        self.text_area.setReadOnly(True)
        layout.addWidget(self.text_area)
        
        # Save buttons
        btn_layout = QHBoxLayout()
        self.save_txt_btn = QPushButton("Save as Text")
        self.save_txt_btn.clicked.connect(lambda: self.save_text('txt'))
        self.save_md_btn = QPushButton("Save as Markdown")
        self.save_md_btn.clicked.connect(lambda: self.save_text('md'))
        self.save_html_btn = QPushButton("Save as HTML")
        self.save_html_btn.clicked.connect(lambda: self.save_text('html'))
        
        btn_layout.addWidget(self.save_txt_btn)
        btn_layout.addWidget(self.save_md_btn)
        btn_layout.addWidget(self.save_html_btn)
        layout.addLayout(btn_layout)
        
        # Disable save buttons initially
        self.save_txt_btn.setEnabled(False)
        self.save_md_btn.setEnabled(False)
        self.save_html_btn.setEnabled(False)
        
    def create_menu_bar(self):
        menubar = self.menuBar()
        
        # File menu
        file_menu = menubar.addMenu("File")
        
        open_action = QAction("Open", self)
        open_action.triggered.connect(self.select_file)
        file_menu.addAction(open_action)
        
        exit_action = QAction("Exit", self)
        exit_action.triggered.connect(self.close)
        file_menu.addAction(exit_action)
        
        # Theme menu
        theme_menu = menubar.addMenu("Theme")
        
        light_action = QAction("Light Mode", self)
        light_action.triggered.connect(lambda: self.set_theme('light'))
        theme_menu.addAction(light_action)
        
        dark_action = QAction("Dark Mode", self)
        dark_action.triggered.connect(lambda: self.set_theme('dark'))
        theme_menu.addAction(dark_action)
        
    def set_theme(self, theme):
        if theme == 'dark':
            self.setStyleSheet("""
                QWidget {
                    background-color: #2d2d2d;
                    color: #ffffff;
                }
                QTextEdit {
                    background-color: #1e1e1e;
                    color: #ffffff;
                }
            """)
        else:
            self.setStyleSheet("")
            
    def select_file(self):
        file_path, _ = QFileDialog.getOpenFileName(
            self,
            "Select PowerPoint File",
            "",
            "PowerPoint Files (*.pptx)"
        )
        
        if file_path:
            self.file_label.setText(file_path)
            self.process_file(file_path)
            
    def process_file(self, file_path):
        self.progress.setVisible(True)
        self.progress.setValue(0)
        
        try:
            text = self.extract_text(file_path)
            self.text_area.setPlainText(text)
            self.save_txt_btn.setEnabled(True)
            self.save_md_btn.setEnabled(True)
            self.save_html_btn.setEnabled(True)
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to process file: {str(e)}")
        finally:
            self.progress.setVisible(False)
            
    def extract_text(self, file_path):
        if self.conversion_type.checkedId() == 0:  # PPTX conversion
            prs = Presentation(file_path)
            text = []
            total_shapes = sum(len(slide.shapes) for slide in prs.slides)
            processed_shapes = 0
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
                    
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    text.append(cell.text)
                    
                    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, "text") and sub_shape.text.strip():
                                text.append(sub_shape.text)
                                
                    # Handle SmartArt shapes
                    if shape.has_chart:
                        chart = shape.chart
                        if hasattr(chart, 'has_title') and chart.has_title:
                            text.append(chart.chart_title.text_frame.text)
                        if hasattr(chart, 'categories'):
                            for category in chart.categories:
                                text.append(category.label)
                        if hasattr(chart, 'series'):
                            for series in chart.series:
                                text.append(series.name)
                                if hasattr(series, 'values'):
                                    for value in series.values:
                                        text.append(str(value))
                    
                    processed_shapes += 1
                    self.progress.setValue(int((processed_shapes / total_shapes) * 100))
            
            return "\n\n".join(text)
        elif self.conversion_type.checkedId() == 1:  # MarkItDown conversion
            md = MarkItDown()
            result = md.convert(file_path)
            return result.text_content  # Access the text content property
        else:  # Spire conversion
            presentation = SpirePresentation()
            presentation.LoadFromFile(file_path)
            extracted_text = []
            
            for slide in presentation.Slides:
                for shape in slide.Shapes:
                    if hasattr(shape, "Nodes"):  # Handle SmartArt shapes
                        for node in shape.Nodes:
                            if hasattr(node, "TextFrame") and node.TextFrame.Text:
                                extracted_text.append(node.TextFrame.Text)
                    elif hasattr(shape, "TextFrame") and shape.TextFrame.Text:
                        extracted_text.append(shape.TextFrame.Text)
            
            presentation.Dispose()
            return "\n\n".join(extracted_text)
    
    def save_text(self, format_type):
        text = self.text_area.toPlainText()
        if not text.strip():
            QMessageBox.warning(self, "Warning", "No text to save")
            return
            
        file_types = {
            'txt': ("Text Files (*.txt)", ".txt"),
            'md': ("Markdown Files (*.md)", ".md"),
            'html': ("HTML Files (*.html)", ".html")
        }
        
        save_path, _ = QFileDialog.getSaveFileName(
            self,
            "Save File",
            "",
            file_types[format_type][0]
        )
        
        if save_path:
            try:
                content = text
                if format_type == 'md':
                    content = markdown.markdown(text)
                elif format_type == 'html':
                    content = f"<html><body><pre>{text}</pre></body></html>"
                
                with open(save_path, 'w', encoding='utf-8') as f:
                    f.write(content)
                QMessageBox.information(self, "Success", "File saved successfully")
            except Exception as e:
                QMessageBox.critical(self, "Error", f"Failed to save file: {str(e)}")

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = PPTXExtractorApp()
    window.show()
    sys.exit(app.exec())
