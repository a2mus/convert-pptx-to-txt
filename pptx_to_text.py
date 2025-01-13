import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext
from markitdown import MarkItDown

from pptx import Presentation

def extract_text_from_pptx(file_path):
    """Extract text from all shapes in a PPTX file"""
    prs = Presentation(file_path)
    text = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            # Extract text from text boxes
            if hasattr(shape, "text"):
                if shape.text.strip():
                    text.append(shape.text)
            
            # Extract text from tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text.append(cell.text)
            
            # Extract text from group shapes
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                for sub_shape in shape.shapes:
                    if hasattr(sub_shape, "text") and sub_shape.text.strip():
                        text.append(sub_shape.text)
    
    return "\n\n".join(text)

class PPTXExtractorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PPTX Text Extractor")
        self.root.geometry("600x400")
        
        # Create GUI elements
        self.create_widgets()
        
    def create_widgets(self):
        # File selection frame
        file_frame = tk.Frame(self.root)
        file_frame.pack(pady=10)
        
        self.file_label = tk.Label(file_frame, text="No file selected")
        self.file_label.pack(side=tk.LEFT, padx=5)
        
        browse_btn = tk.Button(file_frame, text="Browse", command=self.select_file)
        browse_btn.pack(side=tk.LEFT)
        
        # Text display area
        self.text_area = scrolledtext.ScrolledText(self.root, wrap=tk.WORD)
        self.text_area.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)
        
        # Save button
        save_btn = tk.Button(self.root, text="Save Text", command=self.save_text)
        save_btn.pack(pady=10)
        
    def select_file(self):
        file_path = filedialog.askopenfilename(
            filetypes=[("PowerPoint files", "*.pptx")]
        )
        if file_path:
            self.file_label.config(text=file_path)
            try:
                text = extract_text_from_pptx(file_path)
                self.text_area.delete(1.0, tk.END)
                self.text_area.insert(tk.END, text)
            except Exception as e:
                messagebox.showerror("Error", f"Failed to process file: {str(e)}")
    
    def save_text(self):
        text = self.text_area.get(1.0, tk.END)
        if not text.strip():
            messagebox.showwarning("Warning", "No text to save")
            return
            
        save_path = filedialog.asksaveasfilename(
            defaultextension=".txt",
            filetypes=[("Text files", "*.txt")]
        )
        if save_path:
            try:
                with open(save_path, 'w', encoding='utf-8') as f:
                    f.write(text)
                messagebox.showinfo("Success", "Text saved successfully")
            except Exception as e:
                messagebox.showerror("Error", f"Failed to save file: {str(e)}")

if __name__ == "__main__":
    root = tk.Tk()
    app = PPTXExtractorApp(root)
    root.mainloop()
